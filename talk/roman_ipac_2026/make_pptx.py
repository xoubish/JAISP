#!/usr/bin/env python
"""Build jaisp_roman_talk.pptx — PowerPoint version of the JAISP Roman talk.

Mirrors the HTML deck: same 15 slides, dark theme, survey color coding
(Rubin teal / Euclid gold / Roman coral), big-label figures from figures/,
and the per-slide talk track in the PowerPoint notes pane (visible in
presenter view). Regenerate with:  python make_pptx.py
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from PIL import Image
from pathlib import Path

HERE = Path(__file__).resolve().parent
FIGS = HERE / "figures"

GROUND = RGBColor(0x07, 0x0B, 0x14)
PANEL  = RGBColor(0x0E, 0x15, 0x24)
LINE   = RGBColor(0x1E, 0x2A, 0x42)
TEXT   = RGBColor(0xE8, 0xED, 0xF7)
MUTED  = RGBColor(0x93, 0xA0, 0xB8)
FAINT  = RGBColor(0x5A, 0x6A, 0x88)
GOLD   = RGBColor(0xE8, 0xB4, 0x4C)
TEAL   = RGBColor(0x5B, 0xC8, 0xD4)
CORAL  = RGBColor(0xE2, 0x60, 0x4A)
PLATE  = RGBColor(0xF5, 0xF4, 0xF0)
FAIL   = RGBColor(0xD9, 0x8A, 0x7E)

SERIF = "Palatino Linotype"
SANS  = "Calibri"
MONO  = "Consolas"

W, H = Inches(13.333), Inches(7.5)
prs = Presentation()
prs.slide_width, prs.slide_height = W, H
BLANK = prs.slide_layouts[6]


def add_slide(section, title, notes, title_color=GOLD, pageno=True):
    s = prs.slides.add_slide(BLANK)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = GROUND
    if section:
        eb = s.shapes.add_textbox(Inches(0.55), Inches(0.28), Inches(9), Inches(0.32))
        p = eb.text_frame.paragraphs[0]
        r = p.add_run(); r.text = section.upper()
        r.font.name = MONO; r.font.size = Pt(11); r.font.color.rgb = MUTED
        if pageno:
            pg = s.shapes.add_textbox(Inches(11.9), Inches(0.28), Inches(1), Inches(0.32))
            p2 = pg.text_frame.paragraphs[0]; p2.alignment = PP_ALIGN.RIGHT
            r2 = p2.add_run(); r2.text = f"{len(prs.slides._sldIdLst)} / 15"
            r2.font.name = MONO; r2.font.size = Pt(11); r2.font.color.rgb = FAINT
        ln = s.shapes.add_shape(1, Inches(0.55), Inches(0.62), Inches(12.25), Emu(9525))
        ln.fill.solid(); ln.fill.fore_color.rgb = LINE; ln.line.fill.background()
    if title:
        tb = s.shapes.add_textbox(Inches(0.52), Inches(0.72), Inches(12.3), Inches(0.75))
        p = tb.text_frame.paragraphs[0]
        r = p.add_run(); r.text = title
        r.font.name = SERIF; r.font.size = Pt(28); r.font.bold = True
        r.font.color.rgb = TEXT
        # color the emphasized tail if marked with |
        if "|" in title:
            tb.text_frame.clear()
            pre, em = title.split("|", 1)
            p = tb.text_frame.paragraphs[0]
            r1 = p.add_run(); r1.text = pre
            r1.font.name = SERIF; r1.font.size = Pt(28); r1.font.bold = True; r1.font.color.rgb = TEXT
            r2 = p.add_run(); r2.text = em
            r2.font.name = SERIF; r2.font.size = Pt(28); r2.font.bold = True
            r2.font.italic = True; r2.font.color.rgb = title_color
    if notes:
        s.notes_slide.notes_text_frame.text = notes
    return s


def bullets(slide, x, y, w, items, size=15, gap=8, color=TEXT, accent=GOLD):
    tb = slide.shapes.add_textbox(x, y, w, min(Inches(4), Inches(7.42) - y))
    tf = tb.text_frame; tf.word_wrap = True
    first = True
    for item in items:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.space_after = Pt(gap)
        # segments: list of (text, color, bold)
        for txt, col, bold in item:
            r = p.add_run(); r.text = txt
            r.font.name = SANS; r.font.size = Pt(size)
            r.font.color.rgb = col; r.font.bold = bold
    return tb


def plate(slide, img, x, y, w, caption=None):
    im = Image.open(img)
    h_img = w * im.height / im.width
    pad = Inches(0.09)
    box = slide.shapes.add_shape(1, x - pad, y - pad, w + 2 * pad, h_img + 2 * pad)
    box.fill.solid(); box.fill.fore_color.rgb = PLATE; box.line.fill.background()
    box.shadow.inherit = False
    slide.shapes.add_picture(str(img), x, y, width=w)
    if caption:
        cb = slide.shapes.add_textbox(x, y + h_img + Inches(0.12), w, Inches(0.3))
        p = cb.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = caption.upper()
        r.font.name = MONO; r.font.size = Pt(9.5); r.font.color.rgb = FAINT
    return y + h_img


def card(slide, x, y, w, h, header, header_color=GOLD):
    box = slide.shapes.add_shape(5, x, y, w, h)  # rounded rect
    box.fill.solid(); box.fill.fore_color.rgb = PANEL
    box.line.color.rgb = LINE; box.line.width = Pt(1)
    box.shadow.inherit = False
    hb = slide.shapes.add_textbox(x + Inches(0.22), y + Inches(0.14), w - Inches(0.44), Inches(0.32))
    p = hb.text_frame.paragraphs[0]
    r = p.add_run(); r.text = header.upper()
    r.font.name = MONO; r.font.size = Pt(11); r.font.color.rgb = header_color
    return x + Inches(0.22), y + Inches(0.52), w - Inches(0.44)


T, M, K, R, C = TEXT, MUTED, GOLD, TEAL, CORAL

# ============================================================ 1 TITLE
s = add_slide(None, None, "0:00-0:30 - JAISP: one self-supervised representation of both surveys' pixels, reused for detection and astrometry. Everything shown is from the earliest public data of each survey - and the design extends naturally to Roman as a third instrument.")
tb = s.shapes.add_textbox(Inches(0.9), Inches(1.1), Inches(11), Inches(1.6))
r = tb.text_frame.paragraphs[0].add_run(); r.text = "JAISP"
r.font.name = SERIF; r.font.size = Pt(88); r.font.bold = True; r.font.color.rgb = TEXT
tb = s.shapes.add_textbox(Inches(0.95), Inches(2.75), Inches(10.5), Inches(1.0))
tf = tb.text_frame; tf.word_wrap = True
r = tf.paragraphs[0].add_run()
r.text = "A joint Rubin–Euclid foundation model for cross-instrument source detection and astrometry"
r.font.name = SANS; r.font.size = Pt(21); r.font.color.rgb = TEXT
p = tf.add_paragraph(); r = p.add_run(); r.text = "— and a path to Roman —"
r.font.name = SERIF; r.font.size = Pt(18); r.font.italic = True; r.font.color.rgb = CORAL
chips = [("RUBIN · ugrizy", TEAL), ("EUCLID · VIS YJH", GOLD), ("ROMAN · NEXT", CORAL)]
cx = Inches(0.95)
for txt, col in chips:
    wch = Inches(2.35)
    box = s.shapes.add_shape(5, cx, Inches(4.15), wch, Inches(0.42))
    box.fill.background(); box.line.color.rgb = col; box.line.width = Pt(1.2)
    box.shadow.inherit = False
    p = box.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = txt
    r.font.name = MONO; r.font.size = Pt(11); r.font.color.rgb = col
    cx += wch + Inches(0.25)
tb = s.shapes.add_textbox(Inches(0.95), Inches(5.05), Inches(11.5), Inches(1.3))
tf = tb.text_frame; tf.word_wrap = True
r = tf.paragraphs[0].add_run(); r.text = "Shoubaneh Hemmati · IPAC, Caltech"
r.font.name = SANS; r.font.size = Pt(15); r.font.color.rgb = TEXT
p = tf.add_paragraph()
r = p.add_run(); r.text = "with G. Helou, Y.-H. Lin, R. Akeson, C. D. Dowell, A. Faisst, T. Greene, D. Shupe  ·  Hemmati et al. 2026  ·  github.com/xoubish/JAISP"
r.font.name = SANS; r.font.size = Pt(12); r.font.color.rgb = MUTED
p = tf.add_paragraph(); p.space_before = Pt(16)
r = p.add_run(); r.text = "ROMAN SCIENCE CONFERENCE · CALTECH/IPAC · JULY 2026"
r.font.name = MONO; r.font.size = Pt(11); r.font.color.rgb = FAINT

# ============================================================ 2 WHY
s = add_slide("Motivation", "Why combine surveys at the |pixel level?",
  "0:30-2:00 - Cosmology is systematics-limited. The value of three surveys on the same sky from 0.3-2 um is mutual calibration, and it only fully pays off at the pixel level. Position carefully: not replacing classical pipelines - moving the joint, cross-instrument part of the problem into a learned representation validated against classical controls.")
bullets(s, Inches(0.85), Inches(1.85), Inches(11.6), [
 [("Cosmology is now limited by ", T, False), ("systematics, not statistics", K, True),
  (": the remaining disagreements between cosmological probes are too significant to ignore, but we cannot yet tell measurement error apart from new physics.", T, False)],
 [("Rubin/LSST, Euclid, and Roman observe overlapping sky from ~0.3 to 2 μm at sub-arcsecond resolution — in combination they can ", T, False),
  ("calibrate and cross-check one another", K, True),
  (", suppressing survey-specific systematics no single survey controls alone.", T, False)],
 [("Working at the pixel level means handling two instruments with ", T, False),
  ("different resolutions, PSFs, depths, and noise", R, True),
  (" at the same time — something no single-survey pipeline was built to do.", T, False)],
 [("The goal is not to discard classical pipelines, but to move the joint, cross-instrument part of the problem into a representation ", T, False),
  ("learned from the pixels themselves", K, True), (".", T, False)],
], size=17, gap=14)

# ============================================================ 3 RECIPE
s = add_slide("Approach", "Train one model |once| — reuse it for every measurement".replace("|once|","|once —").replace("— —","—"),
  "2:00-3:00 - Masked-band training: hide one of ten bands, predict it from the other nine. Supervision comes from pixels - no labels, no simulations. Land the cost box: 2.5 GPU-hours, ~9M parameters; every capability is a small add-on network on stored features. Deliberately a SMALL foundation model - the paradigm, not scale, does the work.")
bullets(s, Inches(0.85), Inches(1.9), Inches(7.0), [
 [("The model trains itself: ", T, True),
  ("hide one of the ten bands and ask it to predict that band's image from the other nine. The images are their own supervision — ", T, False),
  ("no labels, no catalogs, no simulations", K, True), (".", T, False)],
 [("Then lock it and reuse it: ", T, True),
  ("every measurement afterwards — detection, positions, and more — is a small add-on network reading the ", T, False),
  ("same stored features", R, True), (", at a tiny fraction of the training cost.", T, False)],
 [("Detection and astrometry are working today; photometry and shape measurement are in preparation — all reading the same features.", T, False)],
], size=16, gap=14)
cx, cy, cw = card(s, Inches(8.3), Inches(1.85), Inches(4.35), Inches(4.3), "Cost structure")
tb = s.shapes.add_textbox(cx, cy, cw, Inches(3.6)); tf = tb.text_frame; tf.word_wrap = True
r = tf.paragraphs[0].add_run(); r.text = "2.5 GPU-hours"
r.font.name = SERIF; r.font.size = Pt(30); r.font.color.rgb = GOLD
p = tf.add_paragraph(); r = p.add_run(); r.text = "total pretraining cost"
r.font.name = SANS; r.font.size = Pt(11); r.font.color.rgb = MUTED
for k, v in [("Model size", "~9M parameters"), ("Features", "computed once per dataset"),
             ("Each new task", "a small add-on network"), ("Each new field", "zero retraining")]:
    p = tf.add_paragraph(); p.space_before = Pt(9)
    r = p.add_run(); r.text = k + ":  "; r.font.name = SANS; r.font.size = Pt(13); r.font.color.rgb = MUTED
    r = p.add_run(); r.text = v; r.font.name = SANS; r.font.size = Pt(13)
    r.font.color.rgb = GOLD if k == "Each new field" else TEXT
    r.font.bold = (k == "Each new field")

# ============================================================ 4 DATA
s = add_slide("Data", "The |earliest| public release of each survey".replace("|earliest|","|earliest|"),
  "3:00-4:00 - Rubin DP1 commissioning (PSF ~1 arcsec, so VIS is the sharp reference throughout) + Euclid Q1. 790 matched ten-band tiles over ~0.2 deg2. Two stresses: (1) all public, FIRST release of each survey - a floor, not a ceiling; (2) EDF-S held out entirely for the transfer test.")
yb = plate(s, FIGS/"fig1_data.png", Inches(3.05), Inches(1.75), Inches(7.2),
  "ECDFS training field · held-out EDF-S · the ten bands of one matched tile")
bullets(s, Inches(0.85), yb + Inches(0.5), Inches(11.6), [
 [("Rubin DP1", R, True), (" (commissioning camera, ugrizy, 0.2″/px, ~1″ seeing) + ", T, False),
  ("Euclid Q1", K, True),
  (" (VIS and near-infrared YJH, 0.1″/px): 790 overlapping ten-band image tiles covering ~0.2 deg² of one deep field — the first public release of each survey.", T, False)],
 [("A second deep field (EDF-S) is set aside untouched", T, True),
  (" — the model never sees it — to test at the end whether everything transfers to new sky.", T, False)],
], size=13.5, gap=6)

# ============================================================ 5 ARCH
s = add_slide("Architecture", "Each instrument at its |native resolution",
  "4:00-4:45 - One branch per instrument at native sampling, merged on a common 0.4 arcsec grid where all ten bands exchange information. Q&A seed: 'why not one shared grid?' - resampling throws away the sharper instrument's information. The why-reconstruction story is the NEXT slide.")
yb = plate(s, FIGS/"fig2_architecture.png", Inches(3.35), Inches(1.75), Inches(6.6),
  "each band through its own small network · one branch per instrument · merged ten-band layer · measurement modules")
bullets(s, Inches(0.85), yb + Inches(0.55), Inches(11.6), [
 [("Rubin stays at 0.2″/px, Euclid at 0.1″/px — ", T, False),
  ("no resampling onto a shared grid", K, True),
  (", so neither instrument's spatial information is thrown away. The branches merge on a common 0.4″ grid, where all ten bands exchange information.", T, False)],
], size=13.5, gap=6)

# ============================================================ 6 JOURNEY
s = add_slide("Architecture", "We didn't start here: |ten iterations| to this design".replace(" to this design",""),
  "4:45-5:45 - The credibility slide: tenth iteration, not a lucky guess. v1-v5 all tried latent alignment; the decisive kill was a simple supervised CNN beating the best JEPA features on astrometry (38 vs 47 mas). The turn: predicting actual pixels leaves no shortcut that discards position. v7-v8 native grids; v9-v10 refinements. Meta-point: similarity objectives can look great and still be useless for measurement - the negative results are the design knowledge.")
rows = [
 ("v1–v2", "Contrastive patch matching — make matching Rubin/Euclid cutouts look similar in feature space", "FAILED · SKY-DOMINATED", FAIL),
 ("v3", "Object-centric matching — discover individual sources, pair them across instruments", "ABANDONED · NO PRECISION", FAIL),
 ("v4–v5", "Feature matching on full tiles (JEPA) — make the two instruments' features agree at corresponding sky positions", "BEATEN BY A SIMPLE CNN", FAIL),
 ("v6", "The turn: predict actual pixels — hide one band, reconstruct its image; no shortcut that discards position", "PARADIGM SHIFT", GOLD),
 ("v7–v8", "Mixed resolution — one branch per instrument at its native pixel scale, merged at a common physical scale", "WORKS", TEAL),
 ("v9–v10", "Refinements — every band gets an equal vote in training; the loss focuses on real sources", "PRODUCTION", TEAL),
]
y = Inches(1.72)
for ver, what, tag, col in rows:
    tb = s.shapes.add_textbox(Inches(0.85), y, Inches(1.0), Inches(0.6))
    r = tb.text_frame.paragraphs[0].add_run(); r.text = ver
    r.font.name = MONO; r.font.size = Pt(12); r.font.color.rgb = FAINT
    tb = s.shapes.add_textbox(Inches(2.0), y, Inches(7.9), Inches(0.72))
    tf = tb.text_frame; tf.word_wrap = True
    head, _, rest = what.partition(" — ")
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = head; r.font.name = SANS; r.font.size = Pt(13); r.font.bold = True; r.font.color.rgb = TEXT
    r = p.add_run(); r.text = " — " + rest; r.font.name = SANS; r.font.size = Pt(13); r.font.color.rgb = MUTED
    tb = s.shapes.add_textbox(Inches(10.15), y, Inches(2.7), Inches(0.6))
    r = tb.text_frame.paragraphs[0].add_run(); r.text = tag
    r.font.name = MONO; r.font.size = Pt(10); r.font.color.rgb = col
    ln = s.shapes.add_shape(1, Inches(0.85), y + Inches(0.72), Inches(11.9), Emu(9525))
    ln.fill.solid(); ln.fill.fore_color.rgb = LINE; ln.line.fill.background()
    y += Inches(0.80)
tb = s.shapes.add_textbox(Inches(0.85), Inches(6.72), Inches(11.9), Inches(0.5))
p = tb.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
r = p.add_run(); r.text = "Features that look similar aren't features that know where — the year of negative results is the design."
r.font.name = SERIF; r.font.size = Pt(16); r.font.italic = True; r.font.color.rgb = GOLD

# ============================================================ 7 RECON
s = add_slide("Validation · 1 of 2", "Did it learn |cross-instrument| structure?".replace(" structure?",""),
  "5:45-6:25 - Held-out tiles, one band withheld: Rubin r, VIS, and NISP H all come back faithfully; residuals confined to bright cores and faint galaxy wings. Key sentence: a band from ONE instrument is predicted from bands of the OTHER - genuine cross-instrument correspondence, not within-instrument copying.")
plate(s, FIGS/"fig3_reconstruction.png", Inches(0.95), Inches(1.8), Inches(4.6),
  "truth · reconstruction · residual — held-out tiles")
bullets(s, Inches(6.3), Inches(2.3), Inches(6.3), [
 [("Withhold one band on held-out tiles; predict it from the other nine.", T, False)],
 [("Rubin r, Euclid VIS, and NISP H all reconstructed faithfully — residuals confined to bright cores and faint galaxy wings.", T, False)],
 [("A band from ", T, False), ("one instrument, predicted from bands of the other", K, True),
  (": genuine cross-instrument correspondence, not within-instrument copying.", T, False)],
], size=16, gap=14)

# ============================================================ 8 PROBE
s = add_slide("Validation · 2 of 2", "…and is it organized |physically?",
  "6:25-7:00 - Two quick checks the features are usable and physical: a simple linear fit reads each band's brightness back out (R2 0.7-0.95); the model leans on bands the way real galaxy colors predict - VIS the sharp scaffold, Rubin bands lean on spectral neighbors, near-IR bands predict one another. If behind schedule this slide can be 20 seconds.")
yb = plate(s, FIGS/"fig4_probe.png", Inches(2.4), Inches(1.75), Inches(8.5),
  "brightness read back from the features · which bands inform which")
bullets(s, Inches(0.85), yb + Inches(0.5), Inches(11.6), [
 [("A simple linear fit reads each band's brightness straight back out of the features (R² ≈ 0.7–0.95) — the information is there and easy to reach. The model leans on bands the way real galaxy colors would predict: ", T, False),
  ("VIS is the sharp scaffold everything is built on", K, True),
  ("; Rubin bands draw on their spectral neighbors; the three near-infrared bands mostly predict one another.", T, False)],
], size=13.5, gap=6)

# ============================================================ 9 DET OVERLAY
s = add_slide("Task 1 · Detection", "A detection network on the |shared features",
  "7:00-7:45 - First measurement. Trained on lightweight classical labels (3-sigma VIS extraction) - deliberately NOT on the MER catalog (MER-trained gives essentially the same detector anyway). On held-out sky the source lists agree on the vast majority; the learned detector returns one detection per resolved galaxy where classical extraction fragments spiral arms.")
plate(s, FIGS/"fig5_detection.png", Inches(0.95), Inches(1.8), Inches(4.7),
  "held-out VIS tile · four source lists overlaid")
bullets(s, Inches(6.4), Inches(2.2), Inches(6.2), [
 [("A small detection network turns the ten-band features into a ", T, False),
  ("source-evidence map", K, True),
  (" — finding sources is still peak-finding, as in classical detection, but on all ten bands at once.", T, False)],
 [("Trained on ", T, False), ("lightweight classical labels", K, True),
  (" (standard 3σ source extraction on VIS) — no curated catalog needed.", T, False)],
 [("Matches the Euclid Q1 catalog on held-out sky, and returns ", T, False),
  ("one detection per galaxy", K, True),
  (" where classical extraction fragments extended systems.", T, False)],
], size=15, gap=13)

# ============================================================ 10 DET DEPTH
s = add_slide("Task 1 · Detection", "Ten bands buy |real depth",
  "7:45-8:55 - Against MER: 93% completeness, 95% purity (an undercount - a VIS-selected catalog can't vouch for real multi-band sources it lacks). Depth by injection: hide dimmed real sources, count recoveries: 50% depth VIS~25.0, +0.45 mag over the same detector on VIS alone. The control: an identical detector on per-band features gains +0.03 - the depth lives in the combined representation.")
plate(s, FIGS/"fig6_detection_performance.png", Inches(0.85), Inches(1.8), Inches(6.1),
  "injection completeness · catalog recovery · working point")
tb = s.shapes.add_textbox(Inches(7.5), Inches(1.75), Inches(5.2), Inches(1.0))
tf = tb.text_frame
r = tf.paragraphs[0].add_run(); r.text = "+0.45 mag"
r.font.name = SERIF; r.font.size = Pt(34); r.font.color.rgb = GOLD
p = tf.add_paragraph(); r = p.add_run(); r.text = "deeper from ten-band fusion · 50% depth VIS ≈ 25.0"
r.font.name = SANS; r.font.size = Pt(12); r.font.color.rgb = MUTED
rows = [("50% DEPTH (VIS MAG)", "TEN-BAND", "PER-BAND", MUTED, False),
        ("all 10 bands", "25.0", "25.1", TEXT, False),
        ("VIS only", "24.6", "25.1", TEXT, False),
        ("multi-band gain", "+0.45", "+0.03", GOLD, True)]
y = Inches(3.0)
for a, b, c, col, bold in rows:
    tb = s.shapes.add_textbox(Inches(7.5), y, Inches(5.4), Inches(0.34))
    p = tb.text_frame.paragraphs[0]
    r = p.add_run(); r.text = f"{a:<24}{b:>10}{c:>12}"
    r.font.name = MONO; r.font.size = Pt(12); r.font.color.rgb = col; r.font.bold = bold
    y += Inches(0.36)
bullets(s, Inches(7.5), Inches(4.7), Inches(5.3), [
 [("Depth is measured by hiding real sources, dimmed to a chosen magnitude, in the images and counting how many are found — ", T, False), ("no reference catalog needed", K, True), (".", T, False)],
 [("Against the Euclid catalog: ", T, False), ("93% completeness, 95% purity", K, True),
  (" (an undercount — some “false” detections are real sources the VIS-selected catalog missed).", T, False)],
 [("An identical detector reading per-band features gains only +0.03 mag: the depth comes from the ", T, False),
  ("combined ten-band representation", R, True), (" itself.", T, False)],
], size=12, gap=7)

# ============================================================ 11 ASTROMETRY
s = add_slide("Task 2 · Astrometry", "Carrying VIS precision to |every band",
  "8:55-10:05 - Raw cross-survey scatter ~50 mas, mostly photon noise centroiding faint galaxies. The correction network shrinks it to 14-17 mas in every band and re-centers the clouds. Honesty check with hidden sources at known positions: 19 mas at S/N=5 vs ~50 classical - nearly the ~17 mas VIS precision. Mechanism: it doesn't beat VIS, it CARRIES VIS-quality positions into bands photon noise says can't reach them.")
yb = plate(s, FIGS/"fig8_head_correction.png", Inches(3.1), Inches(1.75), Inches(7.1),
  "before correction (dashed) vs after (solid) · position offsets between the two surveys")
bullets(s, Inches(0.85), yb + Inches(0.5), Inches(11.6), [
 [("The same sources measured in both surveys disagree by ~50 mas — mostly photon noise in centroiding faint galaxies. A small correction network reads the model's features at each source and shrinks this to ", T, False),
  ("14–17 mas in every band", K, True), (".", T, False)],
 [("Is it real? We hide sources at ", T, False), ("known positions", K, True),
  (": corrected positions err by only 19 mas at S/N = 5 (classical: ~50 mas) — nearly the ~17 mas precision of the sharp VIS channel itself, now carried into every band.", T, False)],
], size=12.5, gap=6)

# ============================================================ 12 CONCORDANCE
s = add_slide("Task 2 · Astrometry", "How well do two Gaia-tied surveys |actually agree?",
  "10:05-11:00 - The unexpected measurement: with per-source scatter controlled, a coherent 9-10 mas pattern emerges between two solutions BOTH tied to Gaia - ~7 mas common to all ten bands. Two entirely different fitting methods recover the same pattern; the correction absorbs it to ~1.5 mas. Pitch as a product: a few-mas map of cross-survey agreement - QA no single pipeline can produce, exactly what a three-survey era needs.")
plate(s, FIGS/"fig10_concordance_field.png", Inches(0.85), Inches(1.8), Inches(5.9),
  "two independent fits · after correction · shared vs per-band parts")
bullets(s, Inches(7.3), Inches(2.0), Inches(5.4), [
 [("Beneath the per-source scatter sits a coherent ", T, False), ("9–10 mas pattern", K, True),
  (" between the two surveys' astrometric solutions — even though both are tied to Gaia. About ", T, False),
  ("7 mas is shared by all ten bands", K, True), (": one cross-survey systematic.", T, False)],
 [("Two entirely different fitting methods recover the same pattern, so it is not an artifact of either; the per-source correction absorbs it (→ 1.5 mas).", T, False)],
 [("The result is a few-milliarcsecond ", T, False),
  ("map of how well two independently calibrated surveys actually agree", T, True),
  (" — quality assurance no single pipeline can produce.", T, False)],
], size=14, gap=12)

# ============================================================ 13 TRANSFER
s = add_slide("Transfer", "A second deep field — |nothing retrained",
  "11:00-11:40 - The foundation-model claim, tested: same locked model, same measurement networks on EDF-S, nothing retrained. Detection 93.4/93.1 vs 93.3/94.5 on the training field; positions land on the same 12-18 mas floor. Honest caveat in one clause: transfer across SKY, not across instruments - which is the segue.")
yb = plate(s, FIGS/"fig11_ood.png", Inches(2.7), Inches(1.75), Inches(7.9),
  "EDF-S, 72 tiles never seen by the model")
bullets(s, Inches(0.85), yb + Inches(0.5), Inches(11.6), [
 [("On the never-seen field, detection reaches ", T, False),
  ("93.4% completeness and 93.1% purity", K, True),
  (" — vs 93.3% / 94.5% on the training field — and positions land on the same ", T, False),
  ("12–18 mas", K, True), (" floor. Covering more sky means ", T, False),
  ("applying the model, not rebuilding it", T, True), (".", T, False)],
], size=13.5, gap=6)

# ============================================================ 14 ROMAN
s = add_slide("Roman", "Roman|: the natural third stream",
  "11:40-13:10 - PROTECT THIS SLIDE. Three beats: (1) architecturally trivial - one branch per instrument IS the design; training needs only overlapping public images, hours of GPU time. (2) Roman upgrades the model: sharp deep near-IR at 0.11 arcsec/px is the infrared analog of the VIS anchor - today the near-IR side is resampled NISP and we can SEE the ~3 mas residual resampling left. (3) The model upgrades Roman: fusion depth, Roman-quality positions in ground-based bands, three-way agreement maps as cross-mission QA. Close on the launch window.",
  title_color=CORAL)
cols = [
 ("Why the model is ready", [
   [("One encoder branch per instrument is the design — adding Roman means ", T, False),
    ("adding a third branch", C, True), (", not redesigning the model.", T, False)],
   [("Training needs only ", T, False), ("overlapping public images", C, True),
    (" — no labels, no simulations — and hours of GPU time, once WFI imaging overlaps Rubin/Euclid sky.", T, False)]]),
 ("What Roman adds", [
   [("Sharp, deep near-infrared imaging at 0.11″/px", C, True),
    (" — the infrared counterpart of the sharp VIS reference, where today we must rely on resampled NISP.", T, False)],
   [("A wide survey of thousands of deg² ", T, False), ("inside the LSST footprint", C, True),
    (", plus deep time-domain fields.", T, False)]]),
 ("What Roman gains", [
   [("Detection that uses all surveys at once — for Rubin+Euclid that was worth ", T, False),
    ("+0.45 mag", C, True), (" — and Roman-quality positions carried into the ground-based bands.", T, False)],
   [("Three-way agreement maps", C, True),
    (" at the few-mas level — and cross-instrument alignment solved inside the model instead of inherited from preprocessing.", T, False)]]),
]
x = Inches(0.7)
for header, items in cols:
    cx, cy, cw = card(s, x, Inches(1.8), Inches(4.0), Inches(4.35), header, CORAL)
    bullets(s, cx, cy, cw, items, size=12.5, gap=10)
    x += Inches(4.15)
tb = s.shapes.add_textbox(Inches(0.7), Inches(6.5), Inches(11.9), Inches(0.5))
p = tb.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
r = p.add_run(); r.text = "The launch window opens August 30 — the recipe is ready."
r.font.name = SERIF; r.font.size = Pt(17); r.font.italic = True; r.font.color.rgb = CORAL

# ============================================================ 15 TAKEAWAYS
s = add_slide("Takeaways", "One representation — tasks, fields, |instruments",
  "13:10-13:40 - Read the bullets nearly verbatim. Final line: one representation, trained once, reused across tasks, fields - and, next, instruments. Then thank the audience.")
bullets(s, Inches(0.85), Inches(1.85), Inches(11.7), [
 [("A small (~9M-parameter) foundation model, trained ", T, False), ("without any labels", K, True),
  (" on public Rubin + Euclid images in ~2.5 GPU-hours, learns a shared ten-band representation precise enough to measure with.", T, False)],
 [("Detection: ", T, True), ("matches the Euclid catalog (93% / 95%), and ten-band fusion pushes ", T, False),
  ("0.45 mag deeper", K, True), (" than VIS alone.", T, False)],
 [("Astrometry: ", T, True), ("VIS-quality positions in every band — ~50 → ", T, False),
  ("14–17 mas", K, True), (", verified at 19 mas on hidden test sources — plus a few-mas ", T, False),
  ("map of the disagreement", K, True), (" between the two surveys' Gaia-tied solutions.", T, False)],
 [("The whole stack ", T, False), ("transfers unchanged", K, True),
  (" to a second deep field — all on the earliest data both surveys have released: a floor, not a ceiling.", T, False)],
 [("Roman slots in as a ", T, False), ("third encoder stream", C, True),
  (": a sharp near-infrared anchor, and three-way cross-survey calibration.", T, False)],
], size=16, gap=15)
tb = s.shapes.add_textbox(Inches(0.85), Inches(6.75), Inches(11.7), Inches(0.4))
p = tb.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
r = p.add_run(); r.text = "HEMMATI ET AL. 2026 · GITHUB.COM/XOUBISH/JAISP · SHEMMATI@CALTECH.EDU"
r.font.name = MONO; r.font.size = Pt(10); r.font.color.rgb = FAINT

out = HERE / "jaisp_roman_talk.pptx"
prs.save(out)
print("saved", out, f"({out.stat().st_size/1048576:.1f} MB, {len(prs.slides._sldIdLst)} slides)")
